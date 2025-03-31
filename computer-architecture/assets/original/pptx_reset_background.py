import os
import sys

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    # from pptx.enum.shapes import MSO_SHAPE # 实际上不需要显式导入这个来设置背景
    # from pptx.enum.dml import MSO_FILL # fill.solid() 会处理类型
except ImportError:
    print("错误：缺少 'python-pptx' 库。")
    print("请使用 pip 安装： pip install python-pptx")
    sys.exit(1)  # 如果库未安装则退出


def set_blank_background_for_pptx(pptx_file_path):
    """
    打开指定的PPTX文件，并将所有幻灯片的背景设置为空白（白色）。

    Args:
        pptx_file_path (str): PPTX文件的路径。
    """
    try:
        print(f"正在处理文件: {pptx_file_path} ...")
        # 打开演示文稿
        prs = Presentation(pptx_file_path)

        # 定义白色
        white_color = RGBColor(255, 255, 255)

        # 遍历每一张幻灯片
        for i, slide in enumerate(prs.slides):
            # 获取幻灯片背景
            background = slide.background
            # 获取背景填充
            fill = background.fill
            # 设置为纯色填充
            fill.solid()
            # 设置填充颜色为白色
            fill.fore_color.rgb = white_color
            # print(f"  - 已设置幻灯片 {i+1} 的背景为白色") # 可选：取消注释以查看详细进度

        # 保存修改后的演示文稿（覆盖原文件）
        prs.save(pptx_file_path)
        print(f"完成处理并保存: {pptx_file_path}")
        return True

    except Exception as e:
        print(f"处理文件 {pptx_file_path} 时发生错误: {e}")
        print("  - 请检查文件是否有效且未被其他程序打开。")
        return False


def process_all_pptx_in_directory(directory="."):
    """
    查找指定目录（默认为当前目录）下的所有 .pptx 文件，
    并调用 set_blank_background_for_pptx 函数处理它们。
    """
    processed_count = 0
    failed_count = 0
    script_name = os.path.basename(__file__) if '__file__' in globals() else "this_script.py"

    print(f"将在目录 '{os.path.abspath(directory)}' 中查找 .pptx 文件...")
    print("-" * 30)

    for filename in os.listdir(directory):
        # 检查文件扩展名是否为 .pptx (不区分大小写)
        if filename.lower().endswith(".pptx"):
            # 确保不处理脚本自身（以防万一它也叫.pptx结尾，虽然不太可能）
            if filename == script_name:
                continue

            file_path = os.path.join(directory, filename)
            # 确保它是一个文件而不是目录
            if os.path.isfile(file_path):
                if set_blank_background_for_pptx(file_path):
                    processed_count += 1
                else:
                    failed_count += 1
                print("-" * 30)  # 添加分隔线

    print("\n======== 处理完成 ========")
    if processed_count == 0 and failed_count == 0:
        print("未找到任何 .pptx 文件进行处理。")
    else:
        print(f"成功处理文件数: {processed_count}")
        print(f"处理失败文件数: {failed_count}")
    print("==========================")


# --- 主程序入口 ---
if __name__ == "__main__":
    print("警告：此脚本将直接修改当前目录下的 PPTX 文件！")
    print("请确保您已备份重要文件。")
    # confirmation = input("输入 'yes' 继续执行: ") # 可选：添加用户确认步骤
    # if confirmation.lower() != 'yes':
    #    print("操作已取消。")
    #    sys.exit(0)

    # 处理当前目录下的文件
    process_all_pptx_in_directory(os.path.dirname(__file__))
