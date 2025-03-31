# -*- coding: utf-8 -*-
import os
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# --- 配置 ---
TARGET_FONT_NAME = "Arial"  # 目标字体名称
# --- 配置结束 ---


def change_text_font(text_frame, font_name):
    """更改文本框内所有文本的字体"""
    if not text_frame:
        return
    changed_count = 0
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            try:
                run.font.name = font_name
                changed_count += 1
            except Exception as e:
                # 某些特殊情况（如公式）可能无法直接修改字体名称
                print(f"  - Warning: Could not change font for a run: {e}", file=sys.stderr)
    # print(f"    Changed font for {changed_count} run(s) in text_frame.") # Uncomment for detailed logging


def process_pptx_file(filepath, target_font):
    """处理单个PPTX文件，替换字体"""
    print(f"Processing file: {filepath}...")
    try:
        prs = Presentation(filepath)

        slides_processed = 0
        shapes_processed = 0
        tables_processed = 0

        # 1. 遍历幻灯片 (Slides)
        for slide in prs.slides:
            slides_processed += 1
            # print(f"  Processing Slide {slides_processed}...")
            for shape in slide.shapes:
                shapes_processed += 1
                # 2. 处理普通形状中的文本框
                if shape.has_text_frame:
                    # print(f"    Found Text Frame in shape {shape.shape_id} ({shape.name})")
                    change_text_font(shape.text_frame, target_font)

                # 3. 处理表格中的文本
                elif shape.has_table:
                    tables_processed += 1
                    # print(f"    Found Table in shape {shape.shape_id} ({shape.name})")
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            change_text_font(cell.text_frame, target_font)

                # 4. 处理组合形状 (Groups) - 递归处理内部形状
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    # print(f"    Found Group Shape {shape.shape_id} ({shape.name}), processing content...")
                    for s_in_group in shape.shapes:
                        shapes_processed += 1
                        if s_in_group.has_text_frame:
                            # print(f"      Found Text Frame in grouped shape {s_in_group.shape_id} ({s_in_group.name})")
                            change_text_font(s_in_group.text_frame, target_font)
                        elif s_in_group.has_table:
                            # print(f"      Found Table in grouped shape {s_in_group.shape_id} ({s_in_group.name})")
                            tables_processed += 1
                            group_table = s_in_group.table
                            for grp_row in group_table.rows:
                                for grp_cell in grp_row.cells:
                                    change_text_font(grp_cell.text_frame, target_font)

                # 注意：此脚本未显式处理图表(Charts)或智能图形(SmartArt)中的文本字体。
                # 这些元素的文本字体修改可能需要更复杂的处理。

        # 5. 保存修改 (覆盖原文件)
        prs.save(filepath)
        print(f"Finished processing {filepath}. Processed {slides_processed} slides, {shapes_processed} shapes, {tables_processed} tables.")

    except Exception as e:
        print(f"Error processing file {filepath}: {e}", file=sys.stderr)


def main():
    """主函数，查找并处理目录下的PPTX文件"""
    script_dir = os.path.dirname(os.path.abspath(__file__))  # 获取脚本所在的目录
    print(f"Looking for .pptx files in: {script_dir}")
    print(f"Target font: {TARGET_FONT_NAME}")
    print("-" * 30)
    print("IMPORTANT: Ensure you have backed up your original PPTX files before proceeding!")
    print("-" * 30)

    # # 可以在这里取消注释以要求用户确认
    # confirm = input("Continue? (yes/no): ")
    # if confirm.lower() != 'yes':
    #     print("Operation cancelled by user.")
    #     return

    found_files = False
    for filename in os.listdir(script_dir):
        if filename.lower().endswith(".pptx"):
            # 避免处理临时文件 (通常以 ~$ 开头)
            if not filename.startswith("~$"):
                found_files = True
                filepath = os.path.join(script_dir, filename)
                process_pptx_file(filepath, TARGET_FONT_NAME)

    if not found_files:
        print("No .pptx files found in the script directory.")
    else:
        print("-" * 30)
        print("Script finished.")
        print("Remember: For the changes to display correctly, 'Microsoft YaHei' font must be installed on the viewing system.")


if __name__ == "__main__":
    main()
