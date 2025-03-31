import os
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import logging

# --- 配置 ---
# 设置日志记录，以便查看详细信息或错误
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# --- 主逻辑 ---


def remove_images_from_masters(pptx_filepath):
    """
    打开一个PPTX文件，移除其所有幻灯片母版中的图片形状。

    Args:
        pptx_filepath (str): PPTX文件的完整路径。

    Returns:
        bool: 如果成功处理并保存了文件，则返回 True，否则返回 False。
    """
    try:
        logging.info(f"正在处理文件: {os.path.basename(pptx_filepath)}")
        prs = Presentation(pptx_filepath)
        image_removed_count = 0

        # 遍历所有幻灯片母版
        for i, slide_master in enumerate(prs.slide_masters):
            master_name = f"母版 {i+1}"  # 基本命名，实际母版可能有自己的名字
            logging.debug(f"  检查 {master_name}...")

            # 查找母版中的所有形状
            # 注意：我们需要先收集所有要删除的形状ID，然后再删除它们
            # 因为在遍历过程中直接修改集合会导致问题
            shapes_to_remove = []
            for shape in slide_master.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    shapes_to_remove.append(shape.shape_id)
                    logging.info(f"    - 在 {master_name} 中找到要删除的图片: Shape ID {shape.shape_id}")

            # 如果找到了要删除的图片
            if shapes_to_remove:
                logging.info(f"  在 {master_name} 中找到 {len(shapes_to_remove)} 张图片，正在删除...")
                
                # 使用更可靠的方法删除形状
                # 我们需要通过XML操作来删除形状
                for shape_id in shapes_to_remove:
                    try:
                        # 获取形状的XML元素
                        shape_element = None
                        for shape in slide_master.shapes:
                            if shape.shape_id == shape_id:
                                shape_element = shape.element
                                break
                        
                        if shape_element is not None:
                            # 获取父元素并删除
                            parent = shape_element.getparent()
                            if parent is not None:
                                parent.remove(shape_element)
                                image_removed_count += 1
                                logging.debug(f"      > 已删除 Shape ID {shape_id}")
                            else:
                                logging.warning(f"      > 无法获取 Shape ID {shape_id} 的父元素，跳过删除。")
                        else:
                            logging.warning(f"      > 无法找到 Shape ID {shape_id} 的元素，跳过删除。")
                    except Exception as remove_err:
                        logging.error(f"      > 删除 Shape ID {shape_id} 时出错: {remove_err}")
            else:
                logging.debug(f"  在 {master_name} 中未找到可直接删除的图片形状。")

            # 检查幻灯片母版布局中的图片
            for j, layout in enumerate(slide_master.slide_layouts):
                layout_name = f"{master_name} 布局 {j+1}"
                layout_shapes_to_remove = []
                
                for shape in layout.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        layout_shapes_to_remove.append(shape.shape_id)
                        logging.info(f"    - 在 {layout_name} 中找到要删除的图片: Shape ID {shape.shape_id}")
                
                if layout_shapes_to_remove:
                    logging.info(f"  在 {layout_name} 中找到 {len(layout_shapes_to_remove)} 张图片，正在删除...")
                    
                    for shape_id in layout_shapes_to_remove:
                        try:
                            shape_element = None
                            for shape in layout.shapes:
                                if shape.shape_id == shape_id:
                                    shape_element = shape.element
                                    break
                            
                            if shape_element is not None:
                                parent = shape_element.getparent()
                                if parent is not None:
                                    parent.remove(shape_element)
                                    image_removed_count += 1
                                    logging.debug(f"      > 已删除 Shape ID {shape_id}")
                                else:
                                    logging.warning(f"      > 无法获取 Shape ID {shape_id} 的父元素，跳过删除。")
                            else:
                                logging.warning(f"      > 无法找到 Shape ID {shape_id} 的元素，跳过删除。")
                        except Exception as remove_err:
                            logging.error(f"      > 删除 Shape ID {shape_id} 时出错: {remove_err}")

        # 如果删除了任何图片，则保存文件
        if image_removed_count > 0:
            logging.info(f"共删除了 {image_removed_count} 张图片。正在保存文件: {os.path.basename(pptx_filepath)}")
            # 创建备份文件
            backup_path = pptx_filepath + ".bak"
            if not os.path.exists(backup_path):
                import shutil
                shutil.copy2(pptx_filepath, backup_path)
                logging.info(f"已创建备份文件: {os.path.basename(backup_path)}")
            
            prs.save(pptx_filepath)
            logging.info(f"文件已保存: {os.path.basename(pptx_filepath)}")
        else:
            logging.info(f"在 {os.path.basename(pptx_filepath)} 的母版中未找到或删除任何图片，文件未修改。")

        return True

    except Exception as e:
        logging.error(f"处理文件 {os.path.basename(pptx_filepath)} 时发生错误: {e}", exc_info=True)
        return False


# --- 脚本入口 ---
if __name__ == "__main__":
    current_directory = os.path.dirname(__file__)
    logging.info(f"脚本启动，将在目录 '{current_directory}' 中查找 .pptx 文件...")

    processed_files = 0
    failed_files = 0

    # 遍历当前目录下的所有文件
    for filename in os.listdir(current_directory):
        # 确保处理的是 .pptx 文件，并且不是临时文件（如 ~$开头的文件）
        if filename.lower().endswith(".pptx") and not filename.startswith("~$"):
            filepath = os.path.join(current_directory, filename)
            if os.path.isfile(filepath):  # 确保是文件而不是目录
                if remove_images_from_masters(filepath):
                    processed_files += 1
                else:
                    failed_files += 1

    logging.info("-" * 30)
    logging.info(f"脚本执行完毕。")
    logging.info(f"成功处理（或检查过无需修改）的文件数: {processed_files}")
    if failed_files > 0:
        logging.warning(f"处理失败的文件数: {failed_files}")
    logging.info("请检查上面的日志输出以获取详细信息。")
    logging.info("提醒：如有重要文件，请确认修改结果是否符合预期。")
